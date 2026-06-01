"""Render an HTML slide deck to a pixel-perfect PPTX + PDF.

Each `.slide` element (fixed 1280x720) is screenshotted at 2x via headless
Chromium, then placed full-bleed (one picture per 16:9 slide) into a .pptx,
and the same images are assembled into a landscape .pdf.

Usage:
    python scripts/build_deck_artifacts.py <input.html> <output_basename> [--outdir DIR]

Produces <outdir>/<output_basename>.pptx and <outdir>/<output_basename>.pdf
"""
import sys
import argparse
import tempfile
from pathlib import Path

from playwright.sync_api import sync_playwright
from PIL import Image
from pptx import Presentation
from pptx.util import Inches


def render_slides(html_path: Path, tmpdir: Path, scale: int = 2) -> list[Path]:
    pngs: list[Path] = []
    with sync_playwright() as p:
        browser = p.chromium.launch()
        page = browser.new_page(
            viewport={"width": 1280, "height": 720},
            device_scale_factor=scale,
        )
        page.goto(html_path.resolve().as_uri(), wait_until="networkidle")
        # Ensure web fonts (Inter) and remote images are loaded before capture.
        try:
            page.evaluate("document.fonts.ready")
        except Exception:
            pass
        page.wait_for_timeout(1200)
        slides = page.query_selector_all(".slide")
        if not slides:
            raise SystemExit("No .slide elements found in " + str(html_path))
        for i, el in enumerate(slides, 1):
            out = tmpdir / f"slide_{i:02d}.png"
            el.scroll_into_view_if_needed()
            el.screenshot(path=str(out))
            pngs.append(out)
            print(f"  rendered slide {i:02d}")
        browser.close()
    return pngs


def build_pptx(pngs: list[Path], out_path: Path) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]
    for png in pngs:
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(png), 0, 0, width=prs.slide_width, height=prs.slide_height
        )
    prs.save(str(out_path))
    print(f"  wrote {out_path}  ({len(pngs)} slides)")


def build_pdf(pngs: list[Path], out_path: Path, dpi: float = 192.0) -> None:
    imgs = [Image.open(p).convert("RGB") for p in pngs]
    imgs[0].save(
        str(out_path),
        save_all=True,
        append_images=imgs[1:],
        resolution=dpi,
    )
    print(f"  wrote {out_path}  ({len(pngs)} pages)")


def main() -> None:
    ap = argparse.ArgumentParser()
    ap.add_argument("html")
    ap.add_argument("basename")
    ap.add_argument("--outdir", default="deliverables")
    args = ap.parse_args()

    html_path = Path(args.html)
    outdir = Path(args.outdir)
    outdir.mkdir(parents=True, exist_ok=True)

    print(f"Rendering {html_path} ...")
    with tempfile.TemporaryDirectory() as td:
        pngs = render_slides(html_path, Path(td))
        build_pptx(pngs, outdir / f"{args.basename}.pptx")
        build_pdf(pngs, outdir / f"{args.basename}.pdf")
    print("Done.")


if __name__ == "__main__":
    main()
